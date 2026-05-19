"""
PowerPoint (.pptx) to Markdown converter.

Conversion approach:
  - Each slide becomes a DocumentSection.
  - Slide title shape → section heading (## Slide N: Title).
  - Text shapes → body paragraphs (top-to-bottom, left-to-right reading order).
  - Tables → Markdown table syntax.
  - Images → extracted to assets folder, referenced in Markdown.
  - Speaker notes → blockquote at end of slide section.
  - SmartArt / grouped shapes → best-effort text extraction.
  - Progress reported per-slide.

Dependency: python-pptx >= 1.0.0
"""

import hashlib
import os
from typing import Optional, Callable

from .confidence import ConfidenceResult
from .markdown_writer import ConversionOutput, DocumentSection, rows_to_markdown_table
from .logger import ConversionLogger


def convert(
    source_file: str,
    alias: str = "",
    output_root: str = "",
    preserve_images: bool = True,
    use_subfolder: bool = True,
    logger: Optional[ConversionLogger] = None,
    progress_callback: Optional[Callable[[float], None]] = None,
) -> ConversionOutput:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    output = ConversionOutput(source_file=source_file, alias=alias)
    confidence = ConfidenceResult(source_file=source_file)
    output.confidence = confidence
    output.engine_used = "python-pptx"

    def log_info(msg):
        if logger: logger.info(msg)
    def log_warn(msg):
        if logger: logger.warning(msg)
        confidence.add_warning(msg)
    def progress(p):
        if progress_callback: progress_callback(p)

    log_info(f"PPTX converter started | file={os.path.basename(source_file)}")
    progress(0.05)

    try:
        prs = Presentation(source_file)
    except Exception as e:
        log_warn(f"python-pptx failed to open file: {e}")
        confidence.text_extraction = "Failed"
        confidence.overall = "Failed"
        return output

    total_slides = len(prs.slides)
    log_info(f"Opened PPTX | slides={total_slides}")

    if total_slides == 0:
        log_warn("Presentation has no slides.")
        confidence.text_extraction = "Low"
        confidence.derive_overall()
        return output

    # Assets directory for image extraction
    assets_dir = None
    rel_prefix = "assets/"
    if preserve_images and output_root:
        from .markdown_writer import assets_dir_for, assets_rel_prefix_for
        assets_dir = assets_dir_for(source_file, output_root, alias, use_subfolder)
        rel_prefix = assets_rel_prefix_for(source_file, alias, use_subfolder)
        os.makedirs(assets_dir, exist_ok=True)

    img_counter = 0
    table_count = 0
    seen_hashes: set[str] = set()  # track image blob hashes for deduplication

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        prog = 0.1 + (slide_idx / total_slides) * 0.8
        progress(prog)

        parts = []

        # ── Slide title ──────────────────────────────────────
        title_text = _get_slide_title(slide)
        heading = f"## Slide {slide_num}" + (f": {title_text}" if title_text else "")
        parts.append(heading)

        # Add to TOC
        output.add_toc_entry(2, heading.lstrip("# "), slide_num)

        # ── Sort shapes by reading order (top-to-bottom, left-to-right) ──
        shapes = sorted(
            slide.shapes,
            key=lambda s: (s.top or 0, s.left or 0),
        )

        for shape in shapes:
            try:
                # Skip the title placeholder — already handled above
                if shape.has_text_frame and shape == slide.shapes.title:
                    continue

                # Table
                if shape.has_table:
                    md_table = _table_to_markdown(shape.table)
                    if md_table:
                        parts.append(md_table)
                        table_count += 1

                # Image / picture
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if preserve_images and assets_dir:
                        img_counter += 1
                        img_ref = _extract_image(
                            shape, img_counter, slide_num,
                            assets_dir, rel_prefix, output, log_info,
                            seen_hashes=seen_hashes,
                        )
                        if img_ref:
                            parts.append(img_ref)

                # Group shape — recurse for text and images
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_text = _extract_group_text(
                        shape, MSO_SHAPE_TYPE=MSO_SHAPE_TYPE,
                        preserve_images=preserve_images, assets_dir=assets_dir,
                        rel_prefix=rel_prefix, output=output, log_info=log_info,
                        slide_num=slide_num, img_counter_ref=[img_counter],
                        seen_hashes=seen_hashes,
                    )
                    img_counter = group_text[1]
                    if group_text[0].strip():
                        parts.append(group_text[0])

                # Regular text frame
                elif shape.has_text_frame:
                    text = _text_frame_to_md(shape.text_frame)
                    if text.strip():
                        parts.append(text)

            except Exception as e:
                log_warn(f"Slide {slide_num}: could not process shape: {e}")

        # ── Speaker notes ────────────────────────────────────
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"\n> **Speaker Notes:**\n> {notes_text.replace(chr(10), chr(10) + '> ')}")

        # Assemble slide section
        body = "\n\n".join(parts)
        output.add_section(heading=heading, body=body, page_number=slide_num)

    progress(0.95)

    # Confidence
    confidence.text_extraction = "High"
    confidence.table_structure = "High" if table_count > 0 else "N/A"
    confidence.document_order = "High"
    confidence.image_extraction = "High" if preserve_images and img_counter > 0 else "N/A"
    confidence.image_placement = "Medium" if preserve_images and img_counter > 0 else "N/A"
    confidence.ocr_confidence = "N/A"
    confidence.add_note("Engine: python-pptx (PowerPoint structured access)")
    if table_count:
        confidence.add_note(f"Extracted {table_count} table(s) from slides.")
    confidence.derive_overall()

    log_info(f"PPTX conversion complete | slides={total_slides} images={img_counter} tables={table_count}")
    progress(1.0)
    return output


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_slide_title(slide) -> str:
    """Extract the slide title from the title placeholder shape."""
    if slide.shapes.title is not None:
        return slide.shapes.title.text.strip()
    # Fallback: look for a shape named "Title" or the first large text
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            name = getattr(shape, 'name', '').lower()
            if 'title' in name:
                return shape.text.strip()
    return ""


def _text_frame_to_md(text_frame) -> str:
    """Convert a text frame's paragraphs to Markdown."""
    parts = []
    for para in text_frame.paragraphs:
        level = para.level or 0
        text = _paragraph_to_md(para)
        if not text.strip():
            continue

        # Bulleted paragraphs (level > 0 or has bullet)
        if level > 0:
            indent = "  " * (level - 1)
            parts.append(f"{indent}- {text}")
        else:
            parts.append(text)

    return "\n".join(parts)


def _paragraph_to_md(para) -> str:
    """Convert a single paragraph's runs to Markdown with inline formatting."""
    fragments = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        if run.font.bold and run.font.italic:
            text = f"***{text}***"
        elif run.font.bold:
            text = f"**{text}**"
        elif run.font.italic:
            text = f"*{text}*"
        fragments.append(text)
    return "".join(fragments)


def _table_to_markdown(table) -> str:
    """Convert a PPTX table shape to Markdown table syntax."""
    rows_data = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.replace("\n", " ").strip()
            cells.append(cell_text)
        rows_data.append(cells)

    if not rows_data:
        return ""

    headers = rows_data[0]
    body_rows = rows_data[1:] if len(rows_data) > 1 else []

    # Pad rows to consistent column count
    col_count = max(len(r) for r in rows_data)
    headers = _pad(headers, col_count)
    body_rows = [_pad(r, col_count) for r in body_rows]

    return rows_to_markdown_table(headers, body_rows)


def _extract_image(shape, counter: int, slide_num: int,
                   assets_dir: str, rel_prefix: str,
                   output: ConversionOutput, log_info,
                   seen_hashes: Optional[set] = None) -> str:
    """Extract an image from a picture shape, save it, return a Markdown reference."""
    try:
        image = shape.image
        blob = image.blob

        # Deduplicate: skip if we've already saved an identical image blob
        blob_hash = hashlib.md5(blob).hexdigest()
        if seen_hashes is not None:
            if blob_hash in seen_hashes:
                log_info(f"Slide {slide_num}: skipping duplicate image (hash={blob_hash[:8]}…)")
                return ""
            seen_hashes.add(blob_hash)

        ext = image.content_type.split("/")[-1]
        if ext == "jpeg":
            ext = "jpg"
        filename = f"slide_{slide_num:03d}_img_{counter:03d}.{ext}"
        img_path = os.path.join(assets_dir, filename)

        with open(img_path, "wb") as fh:
            fh.write(blob)

        rel_path = f"{rel_prefix}{filename}"
        output.asset_paths.append(rel_path)
        log_info(f"Saved image: {filename} ({len(blob)} bytes)")
        return f"![Image from slide {slide_num}]({rel_path})"

    except Exception:
        return ""


def _extract_group_text(group_shape, **kwargs) -> tuple[str, int]:
    """Recursively extract text and images from a grouped shape.

    Returns (markdown_text, updated_img_counter).
    """
    MSO_SHAPE_TYPE = kwargs.get("MSO_SHAPE_TYPE")
    preserve_images = kwargs.get("preserve_images", False)
    assets_dir = kwargs.get("assets_dir")
    rel_prefix = kwargs.get("rel_prefix", "assets/")
    output = kwargs.get("output")
    log_info = kwargs.get("log_info")
    slide_num = kwargs.get("slide_num", 0)
    img_counter_ref = kwargs.get("img_counter_ref", [0])
    seen_hashes = kwargs.get("seen_hashes")

    parts = []
    for shape in group_shape.shapes:
        if hasattr(shape, 'shapes'):
            inner_text, img_counter_ref[0] = _extract_group_text(shape, **kwargs)
            if inner_text:
                parts.append(inner_text)
        elif (MSO_SHAPE_TYPE is not None
              and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
              and preserve_images and assets_dir):
            img_counter_ref[0] += 1
            img_ref = _extract_image(
                shape, img_counter_ref[0], slide_num,
                assets_dir, rel_prefix, output, log_info,
                seen_hashes=seen_hashes,
            )
            if img_ref:
                parts.append(img_ref)
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
    return "\n".join(p for p in parts if p), img_counter_ref[0]


def _pad(lst: list, length: int) -> list:
    return list(lst) + [""] * max(0, length - len(lst))
